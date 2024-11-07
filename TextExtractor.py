import os
import re
import pdfplumber
import logging
import comtypes.client
import docx
import tempfile
from pptx import Presentation


sensitive_info_patterns = [
        r'\b\d{9,12}\b',
        r'\b\d{10}\b',
        r'\b[\w\.-]+@[\w\.-]+\.\w+\b'
    ]
sensitive_info_patterns = [re.compile(pattern) for pattern in sensitive_info_patterns]

bullet_patterns = [
    r'^\s*[\u2022\u25E6\u26AB]\s+',
    r'^\s*[-•∙◦◎⦿⦾]\s+'
]
bullet_patterns = [re.compile(pattern, re.MULTILINE) for pattern in bullet_patterns]


def extract_text_from_pdf(pdf_path):
    text = ""
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                try:
                    filtered_text = page.extract_text() or ""

                    for pattern in sensitive_info_patterns:
                        filtered_text = pattern.sub('', filtered_text)

                    for pattern in bullet_patterns:
                        filtered_text = pattern.sub('', filtered_text)

                    text += filtered_text

                except Exception as e:
                    logging.warning(f"Không thể đọc trang PDF: {str(e)}")
                    continue

    except Exception as e:
        raise Exception(f"Lỗi khi đọc file PDF: {str(e)}")

    if not text.strip():
        raise Exception("Không trích xuất được văn bản từ file PDF")

    return text


def extract_text_from_pptx(pptx_path):
    try:
        prs = Presentation(pptx_path)
        text_runs = []

        for slide_number, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    filtered_text = shape.text

                    for pattern in sensitive_info_patterns:
                        filtered_text = pattern.sub('', filtered_text)

                    for pattern in bullet_patterns:
                        filtered_text = pattern.sub('', filtered_text)

                    slide_text.append(filtered_text)

            if slide_text:
                text_runs.append(f"=== Slide {slide_number} ===")
                text_runs.extend(slide_text)

        text = '\n'.join(text_runs)

    except Exception as e:
        raise Exception(f"Lỗi khi đọc file PPTX: {str(e)}")

    if not text.strip():
        raise Exception("Không trích xuất được văn bản từ file PPTX")

    return text


def convert_ppt_to_pptx(ppt_path):
    """Convert .ppt to .pptx format"""
    try:
        pptx_path = os.path.splitext(ppt_path)[0] + "_converted.pptx"

        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1

        deck = powerpoint.Presentations.Open(ppt_path)
        deck.SaveAs(pptx_path, 24)
        deck.Close()
        powerpoint.Quit()

        return pptx_path

    except Exception as e:
        raise Exception(f"Lỗi khi chuyển đổi PPT sang PPTX: {str(e)}")


def extract_text_from_docx(docx_path):
    try:
        doc = docx.Document(docx_path)
        full_text = []

        for para in doc.paragraphs:
            filtered_text = para.text

            for pattern in sensitive_info_patterns:
                filtered_text = pattern.sub('', filtered_text)

            for pattern in bullet_patterns:
                filtered_text = pattern.sub('', filtered_text)

            full_text.append(filtered_text)

        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    full_text.append(cell.text)

        return '\n'.join(full_text)

    except Exception as e:
        raise Exception(f"Lỗi khi đọc file DOCX: {str(e)}")


def extract_text(uploaded_file):
    # _, ext = os.path.splitext(file_path)
    # file_path = os.path.abspath(file_path)

    try:
        temp_file_path = save_uploaded_file(uploaded_file)
        _, ext = os.path.splitext(uploaded_file.name)

        if ext.lower() == '.pdf':
            text = extract_text_from_pdf(temp_file_path)
        elif ext.lower() == '.pptx':
            text = extract_text_from_pptx(temp_file_path)
        elif ext.lower() == '.ppt':
            pptx_path = convert_ppt_to_pptx(temp_file_path)
            text = extract_text_from_pptx(pptx_path)
            os.remove(pptx_path)
        elif ext.lower() == '.docx':
            text = extract_text_from_docx(temp_file_path)
        else:
            raise ValueError("File không được hỗ trợ!")

        os.remove(temp_file_path)

        return text

    except Exception as e:
        raise Exception(f"Lỗi khi xử lý file: {str(e)}")


def save_uploaded_file(uploaded_file):
    """Save uploaded file to temporary directory and return the file path"""
    try:
        # Create a temporary file with the same extension
        suffix = os.path.splitext(uploaded_file.name)[1]
        with tempfile.NamedTemporaryFile(delete = False, suffix = suffix) as tmp_file:
            tmp_file.write(uploaded_file.getvalue())
            return tmp_file.name
    except Exception as e:
        raise Exception(f"Lỗi khi lưu file tạm: {str(e)}")
